from __future__ import annotations

from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.util import Emu, Inches

from models import PresentationSchema


def build_pptx(
    schema: PresentationSchema,
    cover_image: Optional[bytes] = None,
    slide_images: Optional[list[Optional[bytes]]] = None,
) -> bytes:
    prs = Presentation()
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # --- Title slide (layout 0) ---
    title_layout = prs.slide_layouts[0]
    slide0 = prs.slides.add_slide(title_layout)
    if slide0.shapes.title is not None:
        slide0.shapes.title.text = schema.title[:120] or "Presentation"
    if len(slide0.placeholders) > 1 and schema.subtitle:
        slide0.placeholders[1].text = schema.subtitle[:200]
    # Cover image — inserted as a full-bleed background, then the title/subtitle
    # placeholders are brought to the front so text stays readable.
    if cover_image:
        try:
            pic = slide0.shapes.add_picture(
                BytesIO(cover_image), 0, 0, width=slide_w, height=slide_h
            )
            # Move the picture to the back so the title/subtitle overlay it.
            spTree = pic._element.getparent()
            spTree.remove(pic._element)
            spTree.insert(2, pic._element)
        except Exception:
            pass

    # --- Content slides (layout 1: Title + Content) ---
    content_layout = prs.slide_layouts[1]
    images = slide_images or []
    # Text width halved when an image is present; otherwise full.
    half_w = Inches(5.5)
    full_w = Inches(9)
    left_text = Inches(0.5)
    top_text = Inches(1.6)
    text_h = Inches(5.2)
    img_left = Inches(6.3)
    img_top = Inches(1.6)
    img_w = Inches(3.4)
    img_h = Inches(5.2)

    for idx, slide_data in enumerate(schema.slides[:20]):
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = (slide_data.title or "")[:120]

        # Find the content placeholder and resize it if we will insert an image.
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break

        this_img = images[idx] if idx < len(images) else None
        if body is not None:
            if this_img:
                body.left = left_text
                body.top = top_text
                body.width = half_w
                body.height = text_h
            if slide_data.bullets:
                tf = body.text_frame
                tf.word_wrap = True
                first = True
                for bullet in slide_data.bullets[:10]:
                    bullet = (bullet or "").strip()
                    if not bullet:
                        continue
                    if first:
                        tf.text = bullet[:300]
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet[:300]
                        p.level = 0

        if this_img:
            try:
                slide.shapes.add_picture(
                    BytesIO(this_img), img_left, img_top, width=img_w, height=img_h
                )
            except Exception:
                pass

        if slide_data.notes:
            slide.notes_slide.notes_text_frame.text = slide_data.notes[:2000]

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
